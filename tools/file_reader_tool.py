# tools/file_reader_tool.py
import os
import logging
from pathlib import Path
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, field
from io import BytesIO
import base64

logger = logging.getLogger(__name__)

BASE_DIR = Path(__file__).resolve().parent.parent
STORAGE_DIR = Path(os.getenv("FILES_DIR", BASE_DIR / "storage"))
EXAMPLES_DIR = STORAGE_DIR / "examples"

EXAMPLES_DIR.mkdir(parents=True, exist_ok=True)


@dataclass
class ExtractedImage:
    data: bytes
    filename: str
    content_type: str
    width: Optional[int] = None
    height: Optional[int] = None

    def to_base64(self) -> str:
        return base64.b64encode(self.data).decode('utf-8')


@dataclass
class ExtractedTable:
    headers: List[str]
    rows: List[List[Any]]
    sheet_name: Optional[str] = None

    def to_list(self) -> List[List[Any]]:
        return [self.headers] + self.rows


@dataclass
class ExtractedContent:
    filename: str
    file_type: str
    text: str = ""
    tables: List[ExtractedTable] = field(default_factory=list)
    images: List[ExtractedImage] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)

    def has_content(self) -> bool:
        return bool(self.text or self.tables or self.images)


def find_file(filename: str, role: Optional[str] = None) -> Optional[Path]:
    if EXAMPLES_DIR.exists():
        path = EXAMPLES_DIR / filename
        if path.exists():
            logger.info(f"Найден файл в examples: {path}")
            return path

    if role:
        role_dir = STORAGE_DIR / role
        if role_dir.exists():
            path = role_dir / filename
            if path.exists():
                logger.info(f"Найден файл в папке роли: {path}")
                return path

    for folder in STORAGE_DIR.iterdir():
        if folder.is_dir():
            path = folder / filename
            if path.exists():
                logger.info(f"Найден файл: {path}")
                return path

    logger.warning(f"Файл не найден: {filename}")
    return None


def _extract_from_excel(filepath: Path) -> ExtractedContent:
    from openpyxl import load_workbook

    content = ExtractedContent(
        filename=filepath.name,
        file_type="excel"
    )

    try:
        wb = load_workbook(filepath, data_only=True)

        all_text = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]

            headers = []
            rows = []

            for i, row in enumerate(ws.iter_rows(values_only=True)):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if i == 0:
                    headers = row_data
                else:
                    rows.append(row_data)
                all_text.append(" | ".join(row_data))

            if headers or rows:
                content.tables.append(ExtractedTable(
                    headers=headers,
                    rows=rows,
                    sheet_name=sheet_name
                ))

        content.text = "\n".join(all_text)

        wb_full = load_workbook(filepath)
        for sheet_name in wb_full.sheetnames:
            ws = wb_full[sheet_name]
            if hasattr(ws, '_images'):
                for image in ws._images:
                    try:
                        img_data = image._data()
                        content.images.append(ExtractedImage(
                            data=img_data,
                            filename=f"{sheet_name}_image_{len(content.images)}.png",
                            content_type="image/png"
                        ))
                    except Exception as e:
                        logger.warning(f"Ошибка извлечения картинки из Excel: {e}")

        wb.close()
        wb_full.close()

        content.metadata = {
            "sheets": wb.sheetnames,
            "total_tables": len(content.tables),
            "total_images": len(content.images)
        }

    except Exception as e:
        logger.error(f"Ошибка чтения Excel: {e}")
        raise

    return content


def _extract_from_docx(filepath: Path) -> ExtractedContent:
    from docx import Document

    content = ExtractedContent(
        filename=filepath.name,
        file_type="docx"
    )

    try:
        doc = Document(filepath)

        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        content.text = "\n".join(paragraphs)

        for table in doc.tables:
            headers = []
            rows = []

            for i, row in enumerate(table.rows):
                row_data = [cell.text.strip() for cell in row.cells]
                if i == 0:
                    headers = row_data
                else:
                    rows.append(row_data)

            content.tables.append(ExtractedTable(
                headers=headers,
                rows=rows
            ))

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_part = rel.target_part
                    image_data = image_part.blob

                    ext = Path(rel.target_ref).suffix.lower()
                    content_type = {
                        '.png': 'image/png',
                        '.jpg': 'image/jpeg',
                        '.jpeg': 'image/jpeg',
                        '.gif': 'image/gif',
                        '.bmp': 'image/bmp'
                    }.get(ext, 'image/png')

                    content.images.append(ExtractedImage(
                        data=image_data,
                        filename=f"image_{len(content.images)}{ext}",
                        content_type=content_type
                    ))
                except Exception as e:
                    logger.warning(f"Ошибка извлечения картинки из Word: {e}")

        content.metadata = {
            "paragraphs": len(paragraphs),
            "total_tables": len(content.tables),
            "total_images": len(content.images)
        }

    except Exception as e:
        logger.error(f"Ошибка чтения Word: {e}")
        raise

    return content


def _extract_from_pdf(filepath: Path) -> ExtractedContent:
    import pdfplumber

    content = ExtractedContent(
        filename=filepath.name,
        file_type="pdf"
    )

    try:
        with pdfplumber.open(filepath) as pdf:
            all_text = []

            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    all_text.append(text)

                tables = page.extract_tables()
                for table in tables:
                    if table and len(table) > 0:
                        headers = [str(cell) if cell else "" for cell in table[0]]
                        rows = [[str(cell) if cell else "" for cell in row] for row in table[1:]]
                        content.tables.append(ExtractedTable(
                            headers=headers,
                            rows=rows
                        ))

                if hasattr(page, 'images') and page.images:
                    for img in page.images:
                        try:
                            x0, y0, x1, y1 = img['x0'], img['top'], img['x1'], img['bottom']
                            cropped = page.within_bbox((x0, y0, x1, y1))
                            img_obj = cropped.to_image(resolution=150)

                            img_bytes = BytesIO()
                            img_obj.save(img_bytes, format='PNG')

                            content.images.append(ExtractedImage(
                                data=img_bytes.getvalue(),
                                filename=f"page{page_num}_image_{len(content.images)}.png",
                                content_type="image/png",
                                width=int(x1 - x0),
                                height=int(y1 - y0)
                            ))
                        except Exception as e:
                            logger.warning(f"Ошибка извлечения картинки из PDF: {e}")

            content.text = "\n\n".join(all_text)

            content.metadata = {
                "pages": len(pdf.pages),
                "total_tables": len(content.tables),
                "total_images": len(content.images)
            }

    except Exception as e:
        logger.error(f"Ошибка чтения PDF: {e}")
        raise

    return content


def _extract_from_pptx(filepath: Path) -> ExtractedContent:
    from pptx import Presentation

    content = ExtractedContent(
        filename=filepath.name,
        file_type="pptx"
    )

    try:
        prs = Presentation(filepath)

        all_text = []

        for slide_num, slide in enumerate(prs.slides):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                if shape.has_table:
                    table = shape.table
                    headers = []
                    rows = []

                    for i, row in enumerate(table.rows):
                        row_data = [cell.text.strip() for cell in row.cells]
                        if i == 0:
                            headers = row_data
                        else:
                            rows.append(row_data)

                    content.tables.append(ExtractedTable(
                        headers=headers,
                        rows=rows
                    ))

                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        content.images.append(ExtractedImage(
                            data=image.blob,
                            filename=f"slide{slide_num}_image_{len(content.images)}.{image.ext}",
                            content_type=image.content_type
                        ))
                    except Exception as e:
                        logger.warning(f"Ошибка извлечения картинки из PowerPoint: {e}")

            if slide_text:
                all_text.append(f"--- Слайд {slide_num + 1} ---\n" + "\n".join(slide_text))

        content.text = "\n\n".join(all_text)

        content.metadata = {
            "slides": len(prs.slides),
            "total_tables": len(content.tables),
            "total_images": len(content.images)
        }

    except Exception as e:
        logger.error(f"Ошибка чтения PowerPoint: {e}")
        raise

    return content


def extract_content(filepath: Path) -> ExtractedContent:
    suffix = filepath.suffix.lower()

    extractors = {
        '.xlsx': _extract_from_excel,
        '.xls': _extract_from_excel,
        '.docx': _extract_from_docx,
        '.pdf': _extract_from_pdf,
        '.pptx': _extract_from_pptx,
    }

    extractor = extractors.get(suffix)
    if not extractor:
        content = ExtractedContent(
            filename=filepath.name,
            file_type="text"
        )
        try:
            content.text = filepath.read_text(encoding='utf-8')
        except:
            content.text = filepath.read_text(encoding='latin-1')
        return content

    return extractor(filepath)


def read_multiple_files(filenames: List[str], role: Optional[str] = None) -> List[ExtractedContent]:
    results = []

    for filename in filenames:
        filepath = find_file(filename, role)
        if filepath:
            try:
                content = extract_content(filepath)
                results.append(content)
                logger.info(f"Прочитан файл: {filename} ({content.file_type})")
            except Exception as e:
                logger.error(f"Ошибка чтения {filename}: {e}")
        else:
            logger.warning(f"Файл не найден: {filename}")

    return results


def get_example_files() -> List[Dict[str, Any]]:
    if not EXAMPLES_DIR.exists():
        return []

    files = []
    for f in EXAMPLES_DIR.iterdir():
        if f.is_file():
            files.append({
                "name": f.name,
                "type": f.suffix.lower(),
                "size": f.stat().st_size,
                "path": str(f)
            })

    return files